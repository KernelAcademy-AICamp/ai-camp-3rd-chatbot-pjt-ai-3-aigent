#!/usr/bin/env python3
"""
키워드 소싱 레이더 - 기술 발표 자료 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Alias for consistency
RgbColor = RGBColor

def set_slide_background(slide, r, g, b):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RgbColor(r, g, b)

def add_title_slide(prs, title, subtitle):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 15, 23, 42)  # slate-900

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RgbColor(251, 191, 36)  # amber-400
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, subtitle=""):
    """섹션 구분 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 245, 158, 11)  # amber-500

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RgbColor(15, 23, 42)  # slate-900
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(30, 41, 59)  # slate-800
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, highlight_color=None):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 248, 250, 252)  # slate-50

    # 상단 컬러 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(245, 158, 11)  # amber-500
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(15, 23, 42)  # slate-900

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(51, 65, 85)  # slate-700
        p.space_before = Pt(8)
        p.space_after = Pt(4)

        # 볼드 처리 (•로 시작하는 항목)
        if item.startswith("•"):
            p.level = 0
        elif item.startswith("  -"):
            p.level = 1
            p.font.size = Pt(14)

    return slide

def add_architecture_slide(prs):
    """아키텍처 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 248, 250, 252)

    # 상단 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(245, 158, 11)
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "시스템 아키텍처"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(15, 23, 42)

    # 아키텍처 박스들
    boxes = [
        {"name": "Frontend\n(React/Next.js)", "x": 0.5, "y": 1.5, "w": 2.8, "h": 1.2, "color": (59, 130, 246)},  # blue
        {"name": "API Routes\n(Next.js)", "x": 3.6, "y": 1.5, "w": 2.8, "h": 1.2, "color": (16, 185, 129)},  # green
        {"name": "AI/LLM\n(Groq SDK)", "x": 6.7, "y": 1.5, "w": 2.8, "h": 1.2, "color": (139, 92, 246)},  # purple

        {"name": "Charts\n(Recharts)", "x": 0.5, "y": 3.0, "w": 2.8, "h": 1.0, "color": (236, 72, 153)},  # pink
        {"name": "ML Engine\n(시계열 분석)", "x": 3.6, "y": 3.0, "w": 2.8, "h": 1.0, "color": (249, 115, 22)},  # orange
        {"name": "Supabase\n(DB/Auth)", "x": 6.7, "y": 3.0, "w": 2.8, "h": 1.0, "color": (34, 197, 94)},  # emerald

        {"name": "Naver DataLab API", "x": 1.5, "y": 4.5, "w": 3.0, "h": 0.9, "color": (30, 185, 75)},  # naver green
        {"name": "Coupang Scraper", "x": 5.5, "y": 4.5, "w": 3.0, "h": 0.9, "color": (233, 69, 96)},  # coupang red
    ]

    for box in boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box["x"]), Inches(box["y"]),
            Inches(box["w"]), Inches(box["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*box["color"])
        shape.line.fill.background()

        # 텍스트
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = box["name"]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 데이터 흐름 설명
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(1.5))
    tf = flow_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "데이터 흐름: 사용자 요청 → API Routes → AI 분류 → DataLab/Coupang 데이터 수집 → ML 분석 → LLM 응답 생성 → 차트 시각화"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(100, 116, 139)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_tech_stack_slide(prs):
    """기술 스택 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 248, 250, 252)

    # 상단 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(245, 158, 11)
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "기술 스택"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(15, 23, 42)

    # 카테고리별 기술 스택
    categories = [
        {"title": "Frontend", "items": ["Next.js 16.0.7", "React 19.2.0", "Tailwind CSS v4", "Recharts 3.5.1"], "color": (59, 130, 246)},
        {"title": "AI/ML", "items": ["Groq SDK 0.37.0", "arima 0.2.5", "자체 ML 엔진"], "color": (139, 92, 246)},
        {"title": "Backend", "items": ["Next.js API Routes", "Supabase 2.86.0", "TypeScript"], "color": (16, 185, 129)},
        {"title": "External APIs", "items": ["Naver DataLab", "Coupang Scraper", "Naver Shopping Insight"], "color": (249, 115, 22)},
    ]

    x_positions = [0.3, 2.6, 4.9, 7.2]

    for i, cat in enumerate(categories):
        # 카테고리 박스
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.4),
            Inches(2.2), Inches(5.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(255, 255, 255)
        box.line.color.rgb = RgbColor(*cat["color"])
        box.line.width = Pt(2)

        # 카테고리 타이틀
        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.4),
            Inches(2.2), Inches(0.6)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = RgbColor(*cat["color"])
        title_shape.line.fill.background()

        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = cat["title"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 아이템들
        item_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.1), Inches(2.1),
            Inches(2.0), Inches(4.0)
        )
        tf = item_box.text_frame
        tf.word_wrap = True

        for j, item in enumerate(cat["items"]):
            if j > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = RgbColor(51, 65, 85)
            p.space_before = Pt(6)

    return slide

def add_ai_integration_slide(prs):
    """AI 통합 슬라이드"""
    content = [
        "• AI 모델: Groq SDK (LLM 추론)",
        "  - 빠른 응답 속도 (Groq 인프라)",
        "  - OpenAI 호환 API 인터페이스",
        "",
        "• AI 활용 패턴 (3가지):",
        "  - 질문 유형 분류 (trend/strategy/naming/other)",
        "  - DataLab 실행 여부 판단 (스마트 라우팅)",
        "  - 메인 챗봇 응답 생성 (소싱 전문가 역할)",
        "",
        "• 프롬프트 엔지니어링:",
        "  - 역할 정의: 한국어 소싱 전문가 챗봇",
        "  - 응답 구조화: 데이터 요약 → 제품명 제안 → 다음 액션",
        "  - 니치 키워드 설계 원칙: 타깃/용도/형태/가치 기반 조합",
    ]
    return add_content_slide(prs, "AI 통합 상세", content)

def add_ml_engine_slide(prs):
    """ML 엔진 슬라이드"""
    content = [
        "• 시계열 분석 알고리즘 (10+가지 자체 구현):",
        "  - 선형회귀: 전체 추세 파악 (slope, R², 추세 방향)",
        "  - 지수평활법: 노이즈 제거, 데이터 평활화",
        "  - Holt-Winters: 트렌드 + 계절성 예측",
        "  - Mann-Kendall 검정: 통계적 트렌드 유의성",
        "  - STL 분해: 트렌드/계절성/잔차 분리",
        "  - ARIMA: 시계열 예측 (arima 패키지)",
        "",
        "• 종합 점수 산출:",
        "  - growthScore (성장성): 0-100",
        "  - stabilityScore (안정성): 0-100",
        "  - seasonalityScore (계절성 강도): 0-100",
        "  - recommendation: 5단계 추천 등급",
    ]
    return add_content_slide(prs, "ML 시계열 분석 엔진", content)

def add_data_flow_slide(prs):
    """데이터 흐름 슬라이드"""
    content = [
        "• 핵심 분석 파이프라인:",
        "",
        "  1️⃣ 사용자 메시지 수신",
        "  2️⃣ inferDatalabParams: 연도/카테고리 자동 추출",
        "  3️⃣ decideDatalabByLLM: AI가 분석 필요 여부 판단",
        "  4️⃣ fetchTopKeywords: 네이버 Top 10 키워드 수집",
        "  5️⃣ callShoppingCategoryKeywords: 키워드별 시계열 데이터",
        "  6️⃣ analyzeAdvancedTrend: ML 분석 수행",
        "  7️⃣ Supabase Insert: DB 로깅",
        "  8️⃣ LLM 응답 생성: datalabSummary 포함",
        "  9️⃣ 프론트엔드: keywordInsights 시각화",
    ]
    return add_content_slide(prs, "데이터 처리 흐름", content)

def add_visualization_slide(prs):
    """시각화 슬라이드"""
    content = [
        "• Recharts 기반 인터랙티브 차트 (5가지):",
        "",
        "  📊 GrowthScoreComparisonChart",
        "     - 키워드별 성장점수 비교 (수평 막대)",
        "     - Brush 슬라이더, 클릭 선택 기능",
        "",
        "  📈 TimeSeriesForecastChart",
        "     - 실제값/평활/추세/예측선 복합 차트",
        "     - 드래그 줌, 더블클릭 리셋 기능",
        "",
        "  🎯 OverallScoreRadarChart",
        "     - 5축 레이더 차트 (종합 점수)",
        "",
        "  📋 RisingKeywordsSummary / SeasonalityChart",
        "     - 상승추세 요약 카드, 월별 계절성 패턴",
    ]
    return add_content_slide(prs, "데이터 시각화", content)

def add_external_api_slide(prs):
    """외부 API 슬라이드"""
    content = [
        "• 네이버 DataLab API:",
        "  - 쇼핑 카테고리 트렌드 API",
        "  - 카테고리별 키워드 트렌드 API",
        "  - X-Naver-Client-Id/Secret 인증",
        "",
        "• 네이버 쇼핑인사이트 크롤링:",
        "  - Top 20 키워드 조회",
        "  - 429 에러 시 exponential backoff",
        "  - 키워드 정규화 로직",
        "",
        "• 쿠팡 가격 스크래핑:",
        "  - 키워드별 가격 통계 (min/max/avg)",
        "  - 랜덤 User-Agent, Rate limiting 방지",
    ]
    return add_content_slide(prs, "외부 API 통합", content)

def add_key_features_slide(prs):
    """주요 기능 슬라이드"""
    content = [
        "• 핵심 가치:",
        "  - 데이터 기반 의사결정: 네이버 시계열 + ML 분석",
        "  - 대화형 인터페이스: GPT 기반 자연어 상호작용",
        "  - 실시간 시각화: 성장점수, 예측 그래프",
        "  - 니치 마켓 발굴: 경쟁도 낮은 키워드 추천",
        "",
        "• 스마트 기능:",
        "  - 메시지 내 조건 자동 감지 및 분석 모드 전환",
        "  - LLM 기반 DataLab 실행 필요성 판단",
        "  - 16개 네이버 쇼핑 카테고리 매핑",
        "",
        "• 안정성:",
        "  - 429 Rate Limit: Exponential backoff",
        "  - Coupang 403: 조기 종료 및 안내 메시지",
    ]
    return add_content_slide(prs, "주요 기능 및 특징", content)

def add_summary_slide(prs):
    """요약 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 15, 23, 42)  # slate-900

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(251, 191, 36)  # amber-400
    p.alignment = PP_ALIGN.CENTER

    # 요약 내용
    summary_items = [
        "✅ Next.js 16 + React 19 기반 풀스택 아키텍처",
        "✅ Groq SDK를 통한 LLM 통합 (3가지 용도)",
        "✅ 10+ 시계열 분석 알고리즘 자체 구현",
        "✅ 네이버 DataLab + 쿠팡 데이터 통합",
        "✅ Recharts 기반 인터랙티브 시각화",
        "✅ Supabase 인증 및 데이터 저장",
    ]

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(summary_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(226, 232, 240)  # slate-200
        p.space_before = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # 하단 메시지
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI 기반 이커머스 키워드 분석 챗봇"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(148, 163, 184)  # slate-400
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 추가
    add_title_slide(prs, "키워드 소싱 레이더", "Copaung Code Command - 기술 아키텍처 및 AI 스펙")

    add_section_slide(prs, "01. 프로젝트 개요", "AI 기반 이커머스 키워드 분석 챗봇")

    add_content_slide(prs, "프로젝트 소개", [
        "• 프로젝트명: 키워드 소싱 레이더 (Copaung Code Command)",
        "",
        "• 목적:",
        "  - 네이버 데이터랩 기반 검색 트렌드 분석",
        "  - AI를 활용한 '미래 유망 키워드' 예측",
        "  - '니치 제품명' 자동 추천",
        "",
        "• 타깃 사용자:",
        "  - 1인 셀러",
        "  - 소규모 이커머스 브랜드",
        "",
        "• 핵심 가치:",
        "  - 데이터 기반 의사결정 지원",
        "  - 대화형 자연어 인터페이스",
    ])

    add_section_slide(prs, "02. 시스템 아키텍처", "전체 구조 및 기술 스택")

    add_architecture_slide(prs)
    add_tech_stack_slide(prs)

    add_section_slide(prs, "03. AI 통합", "LLM 및 ML 엔진")

    add_ai_integration_slide(prs)
    add_ml_engine_slide(prs)

    add_section_slide(prs, "04. 데이터 파이프라인", "데이터 흐름 및 외부 API")

    add_data_flow_slide(prs)
    add_external_api_slide(prs)

    add_section_slide(prs, "05. 시각화 및 UX", "데이터 시각화 컴포넌트")

    add_visualization_slide(prs)
    add_key_features_slide(prs)

    add_summary_slide(prs)

    # 저장
    output_path = "/Users/admin/Documents/david/_dev/ai-camp-3rd-chatbot-pjt-ai-3-aigent/키워드소싱레이더_기술발표.pptx"
    prs.save(output_path)
    print(f"✅ 프레젠테이션이 생성되었습니다: {output_path}")

if __name__ == "__main__":
    main()
